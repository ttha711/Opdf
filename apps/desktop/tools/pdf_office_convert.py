import json
import sys
from pathlib import Path

import pdfplumber
from pdf2docx import Converter
from pptx import Presentation
from openpyxl import Workbook
from docx import Document


def extract_pages_text(pdf_path: Path):
    pages = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text(layout=True) or ""
            pages.append((idx, text.strip()))
    return pages


def to_docx(pdf_path: Path, out_path: Path):
    cv = Converter(str(pdf_path))
    try:
        cv.convert(str(out_path))
    finally:
        cv.close()


def _looks_like_heading(line: str) -> bool:
    s = line.strip()
    if not s or len(s) > 90:
        return False
    words = s.split()
    if len(words) <= 10 and any(ch.isupper() for ch in s) and s[:1].isupper():
        return True
    return s.endswith(":")


def _looks_like_list_item(line: str) -> bool:
    s = line.strip()
    if not s:
        return False
    return s.startswith(("- ", "* ")) or (len(s) > 2 and s[0].isdigit() and s[1:3] in [". ", ") "])


def to_docx_rewrite(pdf_path: Path, out_path: Path):
    doc = Document()
    doc.add_heading(f"Rewritten from PDF: {pdf_path.name}", level=1)
    with pdfplumber.open(str(pdf_path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text(layout=True) or ""
            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            if idx > 1:
                doc.add_page_break()
            doc.add_heading(f"Page {idx}", level=2)

            buffer = []

            def flush_paragraph():
                nonlocal buffer
                if buffer:
                    paragraph_text = " ".join(buffer).strip()
                    if paragraph_text:
                        doc.add_paragraph(paragraph_text)
                buffer = []

            for line in lines:
                if _looks_like_heading(line):
                    flush_paragraph()
                    doc.add_heading(line.rstrip(":"), level=3)
                    continue
                if _looks_like_list_item(line):
                    flush_paragraph()
                    clean = line.lstrip("-* ").strip()
                    p = doc.add_paragraph(clean)
                    p.style = "List Bullet"
                    continue
                buffer.append(line)
                if len(buffer) >= 4:
                    flush_paragraph()

            flush_paragraph()

    doc.save(str(out_path))


def to_pptx(pdf_path: Path, out_path: Path):
    prs = Presentation()
    for page_num, text in extract_pages_text(pdf_path):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        if title:
            title.text = f"Page {page_num}"
        body = slide.shapes.add_textbox(36, 96, 880, 420).text_frame
        body.word_wrap = True
        body.text = text if text else "(empty)"
    prs.save(str(out_path))


def to_xlsx(pdf_path: Path, out_path: Path):
    wb = Workbook()
    ws = wb.active
    ws.title = "PDF Text"
    ws.append(["Page", "Text"])
    for page_num, text in extract_pages_text(pdf_path):
        ws.append([page_num, text])
    wb.save(str(out_path))


def main():
    if len(sys.argv) < 4:
        print(json.dumps({"ok": False, "error": "Usage: script <input.pdf> <format> <output> [mode]"}))
        sys.exit(2)

    in_path = Path(sys.argv[1]).resolve()
    fmt = sys.argv[2].lower()
    out_path = Path(sys.argv[3]).resolve()
    mode = (sys.argv[4] if len(sys.argv) > 4 else "auto").lower()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        if fmt == "docx":
            if mode in ["rewrite", "ai-rewrite", "editable"]:
                to_docx_rewrite(in_path, out_path)
            else:
                to_docx(in_path, out_path)
        elif fmt == "pptx":
            to_pptx(in_path, out_path)
        elif fmt == "xlsx":
            to_xlsx(in_path, out_path)
        else:
            raise ValueError(f"Unsupported format: {fmt}")
        print(json.dumps({"ok": True, "output": str(out_path), "mode": mode}))
    except Exception as exc:
        print(json.dumps({"ok": False, "error": str(exc)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
